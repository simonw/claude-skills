#!/usr/bin/env python3
"""
Generate WGU Google Skills Strategic Response Presentation
For: Jason Levin, VP WGU Labs
Date: Q4 2025
"""

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN, PP_PARAGRAPH_ALIGNMENT
from pptx.dml.color import RGBColor

# WGU Brand Colors
WGU_BLUE = RGBColor(0, 63, 135)  # WGU Primary Blue
WGU_DARK_BLUE = RGBColor(0, 40, 85)
WGU_LIGHT_BLUE = RGBColor(100, 149, 237)
WHITE = RGBColor(255, 255, 255)
BLACK = RGBColor(0, 0, 0)
GRAY = RGBColor(128, 128, 128)
GREEN = RGBColor(34, 139, 34)
RED = RGBColor(220, 20, 60)

def add_title_slide(prs):
    """Slide 1: Title Slide"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])  # Blank layout

    # Background
    background = slide.background
    fill = background.fill
    fill.solid()
    fill.fore_color.rgb = WGU_BLUE

    # Main Title
    title_box = slide.shapes.add_textbox(Inches(0.5), Inches(2), Inches(9), Inches(1.5))
    title_frame = title_box.text_frame
    title_frame.word_wrap = True
    p = title_frame.paragraphs[0]
    p.text = "Google Skills Launch: WGU's Strategic Response in AI Education"
    p.font.size = Pt(36)
    p.font.bold = True
    p.font.color.rgb = WHITE
    p.alignment = PP_ALIGN.CENTER

    # Subtitle
    subtitle_box = slide.shapes.add_textbox(Inches(0.5), Inches(3.7), Inches(9), Inches(0.6))
    subtitle_frame = subtitle_box.text_frame
    p = subtitle_frame.paragraphs[0]
    p.text = "Turning Google's Skills Feeder into WGU's Degree Pipeline"
    p.font.size = Pt(24)
    p.font.italic = True
    p.font.color.rgb = WGU_LIGHT_BLUE
    p.alignment = PP_ALIGN.CENTER

    # Presenter Info
    info_box = slide.shapes.add_textbox(Inches(0.5), Inches(5), Inches(9), Inches(1))
    info_frame = info_box.text_frame
    info_frame.paragraphs[0].text = "Presenter: Jason Levin, VP, WGU Labs"
    info_frame.paragraphs[0].font.size = Pt(18)
    info_frame.paragraphs[0].font.color.rgb = WHITE
    info_frame.paragraphs[0].alignment = PP_ALIGN.CENTER

    p2 = info_frame.add_paragraph()
    p2.text = "Date: Q4 2025 Board Meeting"
    p2.font.size = Pt(18)
    p2.font.color.rgb = WHITE
    p2.alignment = PP_ALIGN.CENTER

    # Footer
    footer_box = slide.shapes.add_textbox(Inches(0.5), Inches(7), Inches(9), Inches(0.3))
    footer_frame = footer_box.text_frame
    p = footer_frame.paragraphs[0]
    p.text = "Confidential | WGU Labs | Contact: Jason Levin (jlevin@wgu.edu)"
    p.font.size = Pt(10)
    p.font.color.rgb = WGU_LIGHT_BLUE
    p.alignment = PP_ALIGN.CENTER

def add_content_slide(prs, title, content_items, footer_note=""):
    """Helper function to add a content slide"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])  # Blank layout

    # Title bar with WGU blue background
    title_box = slide.shapes.add_textbox(Inches(0), Inches(0), Inches(10), Inches(1))
    title_box.fill.solid()
    title_box.fill.fore_color.rgb = WGU_BLUE
    title_frame = title_box.text_frame
    p = title_frame.paragraphs[0]
    p.text = title
    p.font.size = Pt(28)
    p.font.bold = True
    p.font.color.rgb = WHITE
    p.alignment = PP_ALIGN.CENTER
    title_frame.vertical_anchor = 1  # Middle

    # Content area
    content_box = slide.shapes.add_textbox(Inches(0.5), Inches(1.3), Inches(9), Inches(5.5))
    content_frame = content_box.text_frame
    content_frame.word_wrap = True

    for i, item in enumerate(content_items):
        if i > 0:
            p = content_frame.add_paragraph()
        else:
            p = content_frame.paragraphs[0]

        p.text = item
        p.font.size = Pt(14)
        p.font.color.rgb = BLACK
        p.space_after = Pt(8)
        p.level = 0

        # Bullet point for most items
        if not item.startswith("•"):
            p.text = "• " + item

    # Footer
    footer_box = slide.shapes.add_textbox(Inches(0.5), Inches(7), Inches(9), Inches(0.3))
    footer_frame = footer_box.text_frame
    p = footer_frame.paragraphs[0]
    p.text = "Confidential | WGU Labs | Contact: Jason Levin (jlevin@wgu.edu)"
    p.font.size = Pt(9)
    p.font.color.rgb = GRAY
    p.alignment = PP_ALIGN.CENTER

    if footer_note:
        note_box = slide.shapes.add_textbox(Inches(0.5), Inches(6.5), Inches(9), Inches(0.4))
        note_frame = note_box.text_frame
        p = note_frame.paragraphs[0]
        p.text = f"Note: {footer_note}"
        p.font.size = Pt(11)
        p.font.italic = True
        p.font.color.rgb = WGU_DARK_BLUE
        p.alignment = PP_ALIGN.LEFT

def add_slide_2(prs):
    """Slide 2: Google Skills - The New Reality"""
    content = [
        "Launched: October 21-22, 2025 (official rollout)",
        "Scale: ~3,000 courses/labs/credentials from Google Cloud, DeepMind, Grow with Google, Education",
        "Focus: AI/tech skills; free/low-cost ($0-$49/month); hands-on labs (Gemini-powered)",
        "Impact: 26M+ completions annually; ties to 150+ employer consortium (Deloitte, Walmart)",
        "Market: Google's bet on AI workforce prep, projecting $112B market by 2034",
        "",
        "User Breakdown:",
        "  - 60% beginners (exploration focus)",
        "  - 40% degree-seekers (OUR TARGET)"
    ]
    add_content_slide(prs, "Google Skills: A Unified AI Learning Juggernaut", content,
                     "Google's entry-level certs feed our degree programs – we close the loop.")

def add_slide_3(prs):
    """Slide 3: Not a Competitor - A Feeder"""
    content = [
        "Google's Focus: Micro-credentials (3-6 months, entry roles $50K-$60K)",
        "Track Record: 1M+ global graduates since 2018",
        "",
        "User Breakdown:",
        "  - 60% beginners (no degree needed)",
        "  - 40% seek formal credentials – OUR TARGET MARKET",
        "",
        "Outcomes: 75% report job/promotion/raise within 6 months (entry-level roles)",
        "",
        "The Opportunity: 8.4M seek certificates/associates annually",
        "  → Google's funnel = 20-30% WGU opportunity",
        "",
        "KEY INSIGHT: They accelerate entry; we build careers"
    ]
    add_content_slide(prs, "Google = Skills Starter; WGU = Degree Closer", content)

def add_slide_4(prs):
    """Slide 4: Credit Transfer - Time & Money Saved"""
    content = [
        "ACE Credit Recommendations: Up to 15 credits total per Google certificate",
        "",
        "Google Certificate to WGU Credit Mappings:",
        "  • IT Support Certificate → 12 credits (waives BS IT prerequisites)",
        "    Savings: 1 term ($3,755)",
        "  • Data Analytics Certificate → 12 credits",
        "    Savings: 1 term ($3,755)",
        "  • Project Management Certificate → 9-12 credits",
        "    Savings: $2,800+",
        "",
        "Transfer Acceptance: 75% acceptance rate at schools like WGU",
        "",
        "Impact: Stackable certificates shave 20-30% off WGU time-to-degree",
        "",
        "Technology: Credly badges make transfers instant – no more silos"
    ]
    add_content_slide(prs, "Seamless Google-to-WGU Pathways", content)

def add_slide_5(prs):
    """Slide 5: WGU's Built-In Advantage"""
    content = [
        "WGU Degrees: 120-180 credits; embed 10+ certifications (e.g., CISSP, AWS)",
        "  → Prepares for mid-level roles ($80K+)",
        "",
        "Student Outcomes:",
        "  • +$22K average salary increase (2 years post-graduation)",
        "  • 97% of graduates recommend WGU",
        "  • 87% employed full-time",
        "",
        "Tuition Value:",
        "  • $3,755-$4,125 per 6-month term (unlimited courses)",
        "  • Average $8,010/year vs. $12,660 national average",
        "",
        "The WGU Advantage: Google starts careers; WGU scales them",
        "  → 380,000+ alumni in leadership positions",
        "",
        "Our competency-based learning + certifications = unbeatable hybrid path"
    ]
    add_content_slide(prs, "Why WGU Wins the Long Game", content)

def add_slide_6(prs):
    """Slide 6: Google Bridge Pilot Design"""
    content = [
        "CORE STRATEGY: Co-Branded Transfer Pipeline",
        "  → Free Google cert trials for WGU prospects",
        "  → Auto-ACE credit transfer",
        "  → AI-powered onboarding (Gemini + Aera)",
        "",
        "Pilot Target: 10,000 IT/Healthcare students; Q1-Q2 2026 launch",
        "",
        "Key Features:",
        "  • Gamified learning paths",
        "  • Direct employer consortium access",
        "  • Ethics training (UNESCO-aligned)",
        "  • WGU Labs AI Playground for assessments",
        "",
        "Expected Results:",
        "  • 20% faster degree completions",
        "  • 15% conversion rate from trial to enrollment",
        "",
        "Leverage: Google Accelerator for free tools (50K+ student capacity)"
    ]
    add_content_slide(prs, '"Google Bridge" Pilot: Our Play', content,
                     "Low-risk approach: 7-day Google trials + WGU term structure")

def add_slide_7(prs):
    """Slide 7: Pilot ROI"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])

    # Title
    title_box = slide.shapes.add_textbox(Inches(0), Inches(0), Inches(10), Inches(1))
    title_box.fill.solid()
    title_box.fill.fore_color.rgb = WGU_BLUE
    title_frame = title_box.text_frame
    p = title_frame.paragraphs[0]
    p.text = "Pilot ROI: $67.6M Revenue, $2M Cost"
    p.font.size = Pt(28)
    p.font.bold = True
    p.font.color.rgb = WHITE
    p.alignment = PP_ALIGN.CENTER
    title_frame.vertical_anchor = 1

    # Content
    content_box = slide.shapes.add_textbox(Inches(0.5), Inches(1.3), Inches(9), Inches(5.2))
    content_frame = content_box.text_frame
    content_frame.word_wrap = True

    items = [
        "CONSERVATIVE PROJECTIONS: High-Impact, Low-Risk",
        "",
        "New Students: 9,000/year (15% of Google's 60K U.S. funnel)",
        "",
        "Revenue Calculation:",
        "  9,000 students × $7,510 avg tuition = $67.6M (Year 1)",
        "",
        "Cost: $2M (Google Accelerator + marketing + pilot operations)",
        "",
        "Net ROI: 33:1 return on investment",
        "",
        "Additional Benefits:",
        "  • +20% retention improvement via AI tools",
        "  • 62% exam score boost from adaptive learning",
        "  • $10K salary premium for graduates",
        "  • 90%+ job placement rate",
        "",
        "Funding: Gates Foundation + WGU Labs; scales to $20M+ annual uplift"
    ]

    for i, item in enumerate(items):
        if i > 0:
            p = content_frame.add_paragraph()
        else:
            p = content_frame.paragraphs[0]

        p.text = item
        if item.startswith("  "):
            p.font.size = Pt(13)
            p.level = 1
        elif "CONSERVATIVE" in item or "Revenue Calculation:" in item:
            p.font.size = Pt(16)
            p.font.bold = True
            p.font.color.rgb = WGU_DARK_BLUE
        else:
            p.font.size = Pt(14)

        p.font.color.rgb = BLACK
        p.space_after = Pt(6)

    # Footer
    footer_box = slide.shapes.add_textbox(Inches(0.5), Inches(7), Inches(9), Inches(0.3))
    footer_frame = footer_box.text_frame
    p = footer_frame.paragraphs[0]
    p.text = "Confidential | WGU Labs | Contact: Jason Levin (jlevin@wgu.edu)"
    p.font.size = Pt(9)
    p.font.color.rgb = GRAY
    p.alignment = PP_ALIGN.CENTER

def add_slide_8(prs):
    """Slide 8: AI Integration Roadmap"""
    content = [
        "Q4 2025: Foundation",
        "  • Audit certificate mappings across all Google Skills programs",
        "  • Survey 5,000 students (AI confidence benchmark: 41% national avg)",
        "",
        "2026: Implementation",
        "  • Deploy Gemini for personalized assessments",
        "  • Implement Aera for retention (70% grading time savings)",
        "  • Launch multi-agent AI tutoring system",
        "",
        "Expected Impact:",
        "  • 86% of students already use AI tools",
        "  • 62% average score gains via adaptive learning",
        "  • 25% reduction in equity gaps",
        "",
        "Ethics Framework: UNESCO guidelines for human-centered AI",
        "",
        "Market Position: Lead the $112B AI education market by 2034",
        "",
        "Proven Results: WGU Labs hackathons validate multi-agent AI effectiveness"
    ]
    add_content_slide(prs, "WGU's AI Edge: From Pilot to Pioneer", content)

def add_slide_9(prs):
    """Slide 9: Risks & Mitigations"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])

    # Title
    title_box = slide.shapes.add_textbox(Inches(0), Inches(0), Inches(10), Inches(1))
    title_box.fill.solid()
    title_box.fill.fore_color.rgb = WGU_BLUE
    title_frame = title_box.text_frame
    p = title_frame.paragraphs[0]
    p.text = "Risks & Mitigations"
    p.font.size = Pt(28)
    p.font.bold = True
    p.font.color.rgb = WHITE
    p.alignment = PP_ALIGN.CENTER
    title_frame.vertical_anchor = 1

    # Risks column (left)
    risks_box = slide.shapes.add_textbox(Inches(0.5), Inches(1.5), Inches(4.5), Inches(4.5))
    risks_frame = risks_box.text_frame
    risks_frame.word_wrap = True

    # Risks header
    p = risks_frame.paragraphs[0]
    p.text = "RISKS"
    p.font.size = Pt(20)
    p.font.bold = True
    p.font.color.rgb = RED
    p.alignment = PP_ALIGN.CENTER
    p.space_after = Pt(10)

    risks = [
        "Pipeline Erosion: 20-30% funnel loss if no bridge program",
        "",
        "AI Cheating/Ethics: 66% of students use AI; integrity concerns",
        "",
        "Equity Gaps: Digital divide affects underserved populations (WGU serves 74%)",
        "",
        "Pricing Pressure: Google's $49/month creates cost competition"
    ]

    for risk in risks:
        p = risks_frame.add_paragraph()
        p.text = risk if risk == "" else "• " + risk
        p.font.size = Pt(12)
        p.font.color.rgb = BLACK
        p.space_after = Pt(6)

    # Mitigations column (right)
    mitigations_box = slide.shapes.add_textbox(Inches(5.25), Inches(1.5), Inches(4.5), Inches(4.5))
    mitigations_frame = mitigations_box.text_frame
    mitigations_frame.word_wrap = True

    # Mitigations header
    p = mitigations_frame.paragraphs[0]
    p.text = "MITIGATIONS"
    p.font.size = Pt(20)
    p.font.bold = True
    p.font.color.rgb = GREEN
    p.alignment = PP_ALIGN.CENTER
    p.space_after = Pt(10)

    mitigations = [
        "Auto-transfer system + free trials = 15% conversion rate",
        "",
        "UNESCO-aligned AI ethics training + proctoring systems",
        "",
        "AI-powered support closes gaps by 25%; nonprofit mission focus",
        "",
        "ROI value proposition: $22K salary lift vs. entry-level certs"
    ]

    for mitigation in mitigations:
        p = mitigations_frame.add_paragraph()
        p.text = mitigation if mitigation == "" else "• " + mitigation
        p.font.size = Pt(12)
        p.font.color.rgb = BLACK
        p.space_after = Pt(6)

    # Key message
    message_box = slide.shapes.add_textbox(Inches(0.5), Inches(6.2), Inches(9), Inches(0.5))
    message_frame = message_box.text_frame
    p = message_frame.paragraphs[0]
    p.text = "KEY MESSAGE: Proactive = Advantage; Reactive = Risk"
    p.font.size = Pt(16)
    p.font.bold = True
    p.font.color.rgb = WGU_DARK_BLUE
    p.alignment = PP_ALIGN.CENTER

    # Footer
    footer_box = slide.shapes.add_textbox(Inches(0.5), Inches(7), Inches(9), Inches(0.3))
    footer_frame = footer_box.text_frame
    p = footer_frame.paragraphs[0]
    p.text = "Confidential | WGU Labs | Contact: Jason Levin (jlevin@wgu.edu)"
    p.font.size = Pt(9)
    p.font.color.rgb = GRAY
    p.alignment = PP_ALIGN.CENTER

def add_slide_10(prs):
    """Slide 10: Timeline & Ask"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])

    # Title
    title_box = slide.shapes.add_textbox(Inches(0), Inches(0), Inches(10), Inches(1))
    title_box.fill.solid()
    title_box.fill.fore_color.rgb = WGU_BLUE
    title_frame = title_box.text_frame
    p = title_frame.paragraphs[0]
    p.text = "Timeline & Ask"
    p.font.size = Pt(28)
    p.font.bold = True
    p.font.color.rgb = WHITE
    p.alignment = PP_ALIGN.CENTER
    title_frame.vertical_anchor = 1

    # Timeline content
    content_box = slide.shapes.add_textbox(Inches(0.5), Inches(1.3), Inches(9), Inches(3.5))
    content_frame = content_box.text_frame
    content_frame.word_wrap = True

    timeline_items = [
        "NEXT STEPS: Approve & Accelerate",
        "",
        "Q4 2025:",
        "  • Complete certificate audit and credit mapping",
        "  • Survey 5,000 students on AI confidence and needs",
        "  • Apply for Google Accelerator (free tools for 50K+ students)",
        "",
        "Q1-Q2 2026:",
        "  • Launch 'Google Bridge' pilot with 10,000 students",
        "  • Deploy AI-powered onboarding and assessment tools",
        "",
        "Q3 2026 & Beyond:",
        "  • Scale program based on pilot results",
        "  • Co-develop AI-CBL integration with Google",
        "  • Publish outcomes and position WGU as AI-degree leader"
    ]

    for i, item in enumerate(timeline_items):
        if i > 0:
            p = content_frame.add_paragraph()
        else:
            p = content_frame.paragraphs[0]

        p.text = item
        if "NEXT STEPS" in item:
            p.font.size = Pt(18)
            p.font.bold = True
            p.font.color.rgb = WGU_DARK_BLUE
        elif item.startswith("Q"):
            p.font.size = Pt(15)
            p.font.bold = True
            p.font.color.rgb = WGU_BLUE
        elif item.startswith("  "):
            p.font.size = Pt(13)
            p.level = 1
        else:
            p.font.size = Pt(14)

        p.font.color.rgb = BLACK
        p.space_after = Pt(6)

    # The Ask (highlighted box)
    ask_box = slide.shapes.add_textbox(Inches(1), Inches(5), Inches(8), Inches(1.3))
    ask_box.fill.solid()
    ask_box.fill.fore_color.rgb = WGU_LIGHT_BLUE
    ask_box.line.color.rgb = WGU_DARK_BLUE
    ask_box.line.width = Pt(2)

    ask_frame = ask_box.text_frame
    ask_frame.word_wrap = True
    ask_frame.margin_left = Inches(0.2)
    ask_frame.margin_right = Inches(0.2)
    ask_frame.margin_top = Inches(0.1)
    ask_frame.margin_bottom = Inches(0.1)

    p = ask_frame.paragraphs[0]
    p.text = "BOARD ASK:"
    p.font.size = Pt(16)
    p.font.bold = True
    p.font.color.rgb = WHITE

    p2 = ask_frame.add_paragraph()
    p2.text = "• Approve $2M pilot funding"
    p2.font.size = Pt(14)
    p2.font.color.rgb = WHITE

    p3 = ask_frame.add_paragraph()
    p3.text = "• Position WGU as the AI-degree leader in higher education"
    p3.font.size = Pt(14)
    p3.font.color.rgb = WHITE

    # Closing message
    closing_box = slide.shapes.add_textbox(Inches(0.5), Inches(6.5), Inches(9), Inches(0.4))
    closing_frame = closing_box.text_frame
    p = closing_frame.paragraphs[0]
    p.text = 'From feeder to force multiplier – let\'s own AI education equity. "This aligns with our mission: Affordable outcomes for all."'
    p.font.size = Pt(13)
    p.font.italic = True
    p.font.color.rgb = WGU_DARK_BLUE
    p.alignment = PP_ALIGN.CENTER

    # Footer
    footer_box = slide.shapes.add_textbox(Inches(0.5), Inches(7), Inches(9), Inches(0.3))
    footer_frame = footer_box.text_frame
    p = footer_frame.paragraphs[0]
    p.text = "Confidential | WGU Labs | Contact: Jason Levin (jlevin@wgu.edu)"
    p.font.size = Pt(9)
    p.font.color.rgb = GRAY
    p.alignment = PP_ALIGN.CENTER

def main():
    """Generate the complete presentation"""
    prs = Presentation()
    prs.slide_width = Inches(10)
    prs.slide_height = Inches(7.5)

    print("Creating slide 1: Title...")
    add_title_slide(prs)

    print("Creating slide 2: Google Skills - The New Reality...")
    add_slide_2(prs)

    print("Creating slide 3: Not a Competitor - A Feeder...")
    add_slide_3(prs)

    print("Creating slide 4: Credit Transfer...")
    add_slide_4(prs)

    print("Creating slide 5: WGU's Built-In Advantage...")
    add_slide_5(prs)

    print("Creating slide 6: Google Bridge Pilot...")
    add_slide_6(prs)

    print("Creating slide 7: Pilot ROI...")
    add_slide_7(prs)

    print("Creating slide 8: AI Integration Roadmap...")
    add_slide_8(prs)

    print("Creating slide 9: Risks & Mitigations...")
    add_slide_9(prs)

    print("Creating slide 10: Timeline & Ask...")
    add_slide_10(prs)

    output_file = 'WGU_Google_Skills_Strategic_Response.pptx'
    prs.save(output_file)
    print(f"\n✓ Presentation saved as: {output_file}")
    print(f"  Total slides: {len(prs.slides)}")
    print(f"  Ready for Jason Levin's Q4 2025 Board presentation!")

if __name__ == "__main__":
    main()
